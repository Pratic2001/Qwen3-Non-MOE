"""
extractors.py

Format-specific content extraction. Given raw bytes (or a URL, for
video/audio) plus a content-type/extension hint, dispatch to whichever
extractor can turn that format into clean text.

Every extractor returns the SAME shape, so callers never branch on format:

    {
        "title": str | None,
        "text": str,
        "author": str | None,
        "date": str | None,
        "content_type": str,       # normalized: html, pdf, docx, pptx, xlsx,
                                    # image, video, audio, text, unknown
        "url": str,
        "error": str | None,
        "extra": dict,             # format-specific metadata (page count,
                                    # sheet names, duration, whether OCR/ASR
                                    # was used, etc.) -- useful for QA/audits
    }

Design principles:
- Every heavy dependency (pdfplumber, python-docx, pytesseract, yt-dlp,
  faster-whisper, ...) is imported LAZILY inside the function that needs it,
  so the MCP server still boots and the other extractors still work if one
  optional dependency isn't installed. A missing dependency produces a clear
  {"error": "..."} telling the operator exactly what to `pip install`,
  never a crash.
- Nothing here does network I/O except `extract_video`/`extract_audio`,
  which need yt-dlp to resolve stream URLs and (for the ASR fallback)
  download an audio-only track. Everything else operates on bytes the
  caller already fetched, so server.py stays the single place that owns
  robots.txt / rate-limiting / size caps for GET requests.
"""

from __future__ import annotations

import io
import os
import re
import tempfile
from typing import Optional

# ---------------------------------------------------------------------------
# Content-type detection
# ---------------------------------------------------------------------------

# Maps both file extensions and MIME-type substrings to a normalized kind.
_EXT_MAP = {
    ".html": "html", ".htm": "html", ".xhtml": "html",
    ".pdf": "pdf",
    ".docx": "docx", ".doc": "docx",
    ".pptx": "pptx", ".ppt": "pptx",
    ".xlsx": "xlsx", ".xls": "xlsx", ".csv": "csv",
    ".png": "image", ".jpg": "image", ".jpeg": "image", ".gif": "image",
    ".webp": "image", ".bmp": "image", ".tiff": "image",
    ".mp4": "video", ".webm": "video", ".mkv": "video", ".mov": "video",
    ".avi": "video",
    ".mp3": "audio", ".wav": "audio", ".m4a": "audio", ".ogg": "audio",
    ".flac": "audio",
    ".txt": "text", ".md": "text", ".markdown": "text", ".rst": "text",
    ".json": "text", ".jsonl": "text",
}

_MIME_MAP = {
    "text/html": "html", "application/xhtml+xml": "html",
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml": "docx",
    "application/msword": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml": "xlsx",
    "application/vnd.ms-excel": "xlsx",
    "text/csv": "csv",
    "image/": "image",
    "video/": "video",
    "audio/": "audio",
    "text/plain": "text", "text/markdown": "text",
    "application/json": "text",
}

# Hosts where the actual content is a video/audio stream even though the URL
# has no file extension (YouTube, Vimeo, SoundCloud, ...). yt-dlp knows how
# to resolve these; anything else with no extension is assumed to be HTML.
_KNOWN_VIDEO_AUDIO_HOSTS = {
    "youtube.com", "youtu.be", "vimeo.com", "dailymotion.com", "twitch.tv",
    "soundcloud.com",
}


def detect_content_kind(url: str, content_type: Optional[str] = None) -> str:
    """Best-effort normalized content kind for `url`, preferring an explicit
    HTTP Content-Type header when given, falling back to the URL's file
    extension, then to known streaming hosts, then defaulting to html."""
    if content_type:
        ct = content_type.split(";")[0].strip().lower()
        for prefix, kind in _MIME_MAP.items():
            if ct == prefix or (prefix.endswith("/") and ct.startswith(prefix)):
                return kind

    path = re.sub(r"[?#].*$", "", url).lower()
    _, ext = os.path.splitext(path)
    if ext in _EXT_MAP:
        return _EXT_MAP[ext]

    host = re.sub(r"^www\.", "", re.sub(r"^https?://", "", url).split("/")[0].lower())
    if any(host == h or host.endswith("." + h) for h in _KNOWN_VIDEO_AUDIO_HOSTS):
        return "video"

    return "html"


def _result(text="", title=None, author=None, date=None, content_type="unknown",
            url="", error=None, extra=None) -> dict:
    return {
        "title": title, "text": (text or "").strip(), "author": author,
        "date": date, "content_type": content_type, "url": url,
        "error": error, "extra": extra or {},
    }


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def extract_html(html: str, url: str) -> dict:
    try:
        import trafilatura
        text = trafilatura.extract(html, url=url, favor_recall=True,
                                    include_comments=False, include_tables=True)
        meta = trafilatura.extract_metadata(html)
        title = meta.title if meta else None
        author = meta.author if meta else None
        date = meta.date if meta else None
        if text and len(text.strip()) > 200:
            return _result(text, title, author, date, "html", url)
    except Exception:
        pass

    try:
        from readability import Document
        doc = Document(html)
        title = doc.title()
        summary_html = doc.summary()
        text = re.sub(r"<[^>]+>", " ", summary_html)
        text = re.sub(r"\s+", " ", text).strip()
        if text and len(text) > 200:
            return _result(text, title, None, None, "html", url)
    except Exception:
        pass

    return _result(text="", content_type="html", url=url, error="extraction failed")


# ---------------------------------------------------------------------------
# PDF (text-layer first, OCR fallback for scanned/image-only pages)
# ---------------------------------------------------------------------------

def extract_pdf(data: bytes, url: str = "", ocr_fallback: bool = True,
                 max_pages: int = 300) -> dict:
    try:
        import pdfplumber
    except ImportError:
        return _result(url=url, content_type="pdf",
                        error="pdfplumber not installed (pip install pdfplumber)")

    pages_text = []
    title = None
    n_pages = 0
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            meta = pdf.metadata or {}
            title = meta.get("Title") or None
            n_pages = len(pdf.pages)
            for page in pdf.pages[:max_pages]:
                pages_text.append(page.extract_text() or "")
    except Exception as e:
        return _result(url=url, content_type="pdf", error=f"pdfplumber failed: {e}")

    text = "\n\n".join(t for t in pages_text if t.strip())
    chars_per_page = len(text) / max(1, n_pages)

    # A real text-layer PDF averages hundreds of chars/page. Well under that
    # (scanned pages, image-only PDFs) means the text layer is missing or
    # near-empty -- fall back to OCR if the caller wants it.
    if chars_per_page < 20 and ocr_fallback:
        ocr_result = _ocr_pdf(data, max_pages=min(max_pages, 50))
        if ocr_result["text"]:
            ocr_result["title"] = title
            ocr_result["url"] = url
            ocr_result["extra"]["pages"] = n_pages
            return ocr_result
        # OCR unavailable/failed -- fall through and return whatever thin
        # text layer we found, or the OCR error, rather than silently
        # returning nothing.
        if text.strip():
            return _result(text, title, content_type="pdf", url=url,
                            extra={"pages": n_pages, "ocr_attempted": True,
                                   "ocr_error": ocr_result.get("error")})
        return _result(url=url, content_type="pdf",
                        error=f"no text layer and OCR fallback failed: {ocr_result.get('error')}",
                        extra={"pages": n_pages})

    if not text.strip():
        return _result(url=url, content_type="pdf", error="no extractable text",
                        extra={"pages": n_pages})

    return _result(text, title, content_type="pdf", url=url, extra={"pages": n_pages})


def _ocr_pdf(data: bytes, max_pages: int = 50) -> dict:
    """Rasterize pages and OCR them. Requires poppler (system package,
    `apt install poppler-utils` / `brew install poppler`) for pdf2image, and
    tesseract (`apt install tesseract-ocr`) for pytesseract."""
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
    except ImportError as e:
        return _result(error=f"OCR deps missing ({e}); "
                        "pip install pdf2image pytesseract, plus system "
                        "packages poppler-utils and tesseract-ocr")
    try:
        images = convert_from_bytes(data, dpi=200, first_page=1, last_page=max_pages)
    except Exception as e:
        return _result(error=f"pdf2image rasterization failed (is poppler "
                        f"installed?): {e}")

    texts = []
    for img in images:
        try:
            texts.append(pytesseract.image_to_string(img))
        except Exception as e:
            return _result(error=f"pytesseract OCR failed: {e}")
    text = "\n\n".join(t for t in texts if t.strip())
    if not text.strip():
        return _result(error="OCR produced no text")
    return _result(text, content_type="pdf", extra={"ocr_used": True, "ocr_pages": len(images)})


# ---------------------------------------------------------------------------
# Office documents
# ---------------------------------------------------------------------------

def extract_docx(data: bytes, url: str = "") -> dict:
    try:
        import docx
    except ImportError:
        return _result(url=url, content_type="docx",
                        error="python-docx not installed (pip install python-docx)")
    try:
        doc = docx.Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
        title = None
        try:
            title = doc.core_properties.title or None
        except Exception:
            pass
        if not text.strip():
            return _result(url=url, content_type="docx", error="no text content")
        return _result(text, title, content_type="docx", url=url)
    except Exception as e:
        return _result(url=url, content_type="docx", error=f"python-docx failed: {e}")


def extract_pptx(data: bytes, url: str = "") -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        return _result(url=url, content_type="pptx",
                        error="python-pptx not installed (pip install python-pptx)")
    try:
        prs = Presentation(io.BytesIO(data))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            chunks = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    chunks.append(shape.text_frame.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        chunks.append(" | ".join(c.text.strip() for c in row.cells))
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                chunks.append(f"[notes] {slide.notes_slide.notes_text_frame.text.strip()}")
            if chunks:
                slides_text.append(f"--- Slide {i} ---\n" + "\n".join(chunks))
        text = "\n\n".join(slides_text)
        n_slides = len(prs.slides._sldIdLst)
        if not text.strip():
            return _result(url=url, content_type="pptx", error="no text content",
                            extra={"slides": n_slides})
        return _result(text, content_type="pptx", url=url, extra={"slides": n_slides})
    except Exception as e:
        return _result(url=url, content_type="pptx", error=f"python-pptx failed: {e}")


def extract_xlsx(data: bytes, url: str = "", max_rows_per_sheet: int = 2000) -> dict:
    try:
        import openpyxl
    except ImportError:
        return _result(url=url, content_type="xlsx",
                        error="openpyxl not installed (pip install openpyxl)")
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        sheets_text = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows_per_sheet:
                    break
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            if rows:
                sheets_text.append(f"--- Sheet: {name} ---\n" + "\n".join(rows))
        text = "\n\n".join(sheets_text)
        if not text.strip():
            return _result(url=url, content_type="xlsx", error="no non-empty cells")
        return _result(text, content_type="xlsx", url=url, extra={"sheets": wb.sheetnames})
    except Exception as e:
        return _result(url=url, content_type="xlsx", error=f"openpyxl failed: {e}")


def extract_csv(data: bytes, url: str = "") -> dict:
    try:
        text = data.decode("utf-8", errors="replace")
        if not text.strip():
            return _result(url=url, content_type="csv", error="empty file")
        return _result(text, content_type="csv", url=url)
    except Exception as e:
        return _result(url=url, content_type="csv", error=str(e))


# ---------------------------------------------------------------------------
# Images (OCR)
# ---------------------------------------------------------------------------

def extract_image(data: bytes, url: str = "") -> dict:
    try:
        from PIL import Image
        import pytesseract
    except ImportError as e:
        return _result(url=url, content_type="image",
                        error=f"OCR deps missing ({e}); pip install pillow "
                        "pytesseract, plus system package tesseract-ocr")
    try:
        img = Image.open(io.BytesIO(data))
        text = pytesseract.image_to_string(img)
        text = text.strip()
        if not text or len(text) < 20:
            return _result(url=url, content_type="image",
                            error="no significant text found in image (OCR result too short)")
        return _result(text, content_type="image", url=url,
                        extra={"width": img.width, "height": img.height})
    except Exception as e:
        return _result(url=url, content_type="image", error=f"OCR failed: {e}")


# ---------------------------------------------------------------------------
# Video / audio (captions preferred, ASR transcription fallback)
# ---------------------------------------------------------------------------

class _YtdlpNullLogger:
    """Swallows yt-dlp's debug/warning/error output instead of letting it
    hit stderr. `quiet: True` alone only suppresses yt-dlp's normal
    progress/info output -- warnings (e.g. "No supported JavaScript
    runtime could be found") still print unless a logger is supplied that
    drops them too. Keeps the MCP subprocess's stderr clean so operators
    aren't left wondering whether a WARNING line means the extraction
    actually failed (it usually didn't -- captions/metadata still come
    through fine without a JS runtime; only some signature-gated video/
    audio formats used by the ASR-fallback download path are affected)."""

    def debug(self, msg):
        pass

    def info(self, msg):
        pass

    def warning(self, msg):
        pass

    def error(self, msg):
        pass


def _ytdlp_base_opts(**overrides) -> dict:
    """Shared yt-dlp options for both the caption-fetch and audio-download
    paths: quiet + null logger (see above), plus an optional JS runtime
    (deno/node/bun/etc.) if the operator has one installed and points to
    it via YTDLP_JS_RUNTIME (e.g. `YTDLP_JS_RUNTIME=deno` or a full
    `name:/path/to/binary`). YouTube increasingly requires executing a bit
    of JS to derive signatures for some formats; without a runtime yt-dlp
    still works for captions/metadata (as seen in practice) but may miss
    some audio/video formats used by the ASR-fallback download path -- see
    https://github.com/yt-dlp/yt-dlp/wiki/EJS for install options."""
    opts = {"quiet": True, "no_warnings": True, "logger": _YtdlpNullLogger()}
    js_runtime = os.environ.get("YTDLP_JS_RUNTIME")
    if js_runtime:
        opts["extractor_args"] = {"youtube": {"jsruntime": [js_runtime]}}
    opts.update(overrides)
    return opts


def extract_video(url: str, prefer_captions: bool = True,
                   asr_fallback: bool = True, asr_model_size: str = "base",
                   max_duration_seconds: int = 3600) -> dict:
    """Get a transcript for a video URL (YouTube/Vimeo/direct file/etc.).

    Order of preference:
      1. Existing manual or auto-generated captions, pulled via yt-dlp
         without downloading any audio/video -- fast, free, and usually
         accurate since a human or the platform's own ASR produced them.
      2. If no captions exist and asr_fallback=True, download an audio-only
         stream and transcribe locally with faster-whisper.

    Requires `yt-dlp` for (1) and (2), and `faster-whisper` (+ ffmpeg on
    PATH) for (2) specifically.
    """
    try:
        import yt_dlp
    except ImportError:
        return _result(url=url, content_type="video",
                        error="yt-dlp not installed (pip install yt-dlp)")

    info = None
    try:
        with yt_dlp.YoutubeDL(_ytdlp_base_opts(
                skip_download=True, writesubtitles=prefer_captions,
                writeautomaticsub=prefer_captions, subtitleslangs=["en"])) as ydl:
            info = ydl.extract_info(url, download=False)
    except Exception as e:
        return _result(url=url, content_type="video", error=f"yt-dlp metadata fetch failed: {e}")

    if not info:
        return _result(url=url, content_type="video", error="yt-dlp returned no metadata")

    duration = info.get("duration") or 0
    if duration and duration > max_duration_seconds:
        return _result(url=url, content_type="video",
                        error=f"duration {duration}s exceeds max_duration_seconds="
                        f"{max_duration_seconds}, skipped to bound cost",
                        extra={"duration": duration})

    title = info.get("title")
    date = info.get("upload_date")
    author = info.get("uploader")

    if prefer_captions:
        caption_text = _captions_from_ytdlp_info(info)
        if caption_text and len(caption_text.strip()) > 100:
            return _result(caption_text, title, author, date, "video", url,
                            extra={"source": "captions", "duration": duration})

    if not asr_fallback:
        return _result(url=url, content_type="video",
                        error="no captions available and asr_fallback=False",
                        extra={"duration": duration})

    return _transcribe_via_ytdlp_audio(url, title, author, date, duration, asr_model_size)


def _captions_from_ytdlp_info(info: dict) -> Optional[str]:
    """Download and flatten a caption/subtitle track (VTT/SRT) yt-dlp found,
    preferring manual captions over auto-generated ones."""
    import httpx

    for key in ("subtitles", "automatic_captions"):
        tracks = info.get(key) or {}
        for lang in ("en", "en-US", "en-GB"):
            if lang in tracks and tracks[lang]:
                # Prefer a plain-text/vtt format over json3/ttml.
                fmt = next((f for f in tracks[lang] if f.get("ext") in ("vtt", "srv1", "ttml")),
                           tracks[lang][0])
                try:
                    resp = httpx.get(fmt["url"], timeout=20)
                    return _strip_subtitle_markup(resp.text)
                except Exception:
                    continue
    return None


def _strip_subtitle_markup(raw: str) -> str:
    """Turn VTT/SRT into plain deduplicated prose: drop timestamps, cue
    numbers, and tags, and collapse the repeated-line rolling-caption
    artifact common in auto-generated subtitles."""
    lines = []
    seen_last = None
    for line in raw.splitlines():
        line = line.strip()
        if not line or line.upper().startswith("WEBVTT"):
            continue
        if re.match(r"^\d+$", line):
            continue
        if re.match(r"^[\d:.,]+\s*-->\s*[\d:.,]+", line):
            continue
        line = re.sub(r"<[^>]+>", "", line)
        line = re.sub(r"\[.*?\]", "", line).strip()
        if not line or line == seen_last:
            continue
        lines.append(line)
        seen_last = line
    return " ".join(lines)


def _transcribe_via_ytdlp_audio(url: str, title, author, date, duration,
                                 model_size: str) -> dict:
    try:
        import yt_dlp
    except ImportError:
        return _result(url=url, content_type="video",
                        error="yt-dlp not installed (pip install yt-dlp)")
    try:
        from faster_whisper import WhisperModel
    except ImportError:
        return _result(url=url, content_type="video",
                        error="faster-whisper not installed (pip install "
                        "faster-whisper) -- required for ASR fallback when "
                        "no captions exist. ffmpeg must also be on PATH.")

    with tempfile.TemporaryDirectory() as tmp:
        out_template = os.path.join(tmp, "audio.%(ext)s")
        try:
            with yt_dlp.YoutubeDL(_ytdlp_base_opts(
                    format="bestaudio/best", outtmpl=out_template,
                    postprocessors=[{"key": "FFmpegExtractAudio",
                                      "preferredcodec": "wav"}])) as ydl:
                ydl.download([url])
        except Exception as e:
            return _result(url=url, content_type="video",
                            error=f"yt-dlp audio download failed: {e}",
                            extra={"duration": duration})

        wav_path = os.path.join(tmp, "audio.wav")
        if not os.path.exists(wav_path):
            candidates = [f for f in os.listdir(tmp) if f.startswith("audio.")]
            if not candidates:
                return _result(url=url, content_type="video",
                                error="audio download produced no file")
            wav_path = os.path.join(tmp, candidates[0])

        try:
            model = _get_whisper_model(model_size)
            segments, _info = model.transcribe(wav_path, beam_size=1, vad_filter=True)
            text = " ".join(seg.text.strip() for seg in segments)
        except Exception as e:
            return _result(url=url, content_type="video",
                            error=f"faster-whisper transcription failed: {e}",
                            extra={"duration": duration})

    if not text.strip():
        return _result(url=url, content_type="video", error="ASR produced no text",
                        extra={"duration": duration})
    return _result(text, title, author, date, "video", url,
                    extra={"source": "asr", "asr_model": model_size, "duration": duration})


_whisper_models: dict = {}


def _get_whisper_model(model_size: str):
    """Cache loaded whisper models across calls in this process -- loading
    is the expensive part (seconds to tens of seconds), transcription reuse
    is what makes batch scraping runs feasible."""
    if model_size not in _whisper_models:
        from faster_whisper import WhisperModel
        # int8 compute_type keeps this usable on CPU-only boxes; if a GPU is
        # available, set WHISPER_DEVICE=cuda / WHISPER_COMPUTE_TYPE=float16.
        device = os.environ.get("WHISPER_DEVICE", "cpu")
        compute_type = os.environ.get("WHISPER_COMPUTE_TYPE", "int8" if device == "cpu" else "float16")
        _whisper_models[model_size] = WhisperModel(model_size, device=device, compute_type=compute_type)
    return _whisper_models[model_size]


def extract_audio(url_or_path: str, is_local_file: bool = False, model_size: str = "base") -> dict:
    """Transcribe a direct audio URL or local file (no captions possible for
    raw audio, so this always goes straight to ASR)."""
    try:
        from faster_whisper import WhisperModel  # noqa: F401 (import check)
    except ImportError:
        return _result(url=url_or_path, content_type="audio",
                        error="faster-whisper not installed (pip install faster-whisper)")

    path = url_or_path
    tmp_dir = None
    if not is_local_file:
        import httpx
        tmp_dir = tempfile.mkdtemp()
        path = os.path.join(tmp_dir, "audio_download")
        try:
            with httpx.stream("GET", url_or_path, timeout=60, follow_redirects=True) as resp:
                resp.raise_for_status()
                with open(path, "wb") as f:
                    for chunk in resp.iter_bytes():
                        f.write(chunk)
        except Exception as e:
            return _result(url=url_or_path, content_type="audio", error=f"audio download failed: {e}")

    try:
        model = _get_whisper_model(model_size)
        segments, _info = model.transcribe(path, beam_size=1, vad_filter=True)
        text = " ".join(seg.text.strip() for seg in segments)
    except Exception as e:
        return _result(url=url_or_path, content_type="audio", error=f"transcription failed: {e}")
    finally:
        if tmp_dir:
            import shutil
            shutil.rmtree(tmp_dir, ignore_errors=True)

    if not text.strip():
        return _result(url=url_or_path, content_type="audio", error="ASR produced no text")
    return _result(text, content_type="audio", url=url_or_path,
                    extra={"source": "asr", "asr_model": model_size})


# ---------------------------------------------------------------------------
# Plain text passthrough
# ---------------------------------------------------------------------------

def extract_text(data: bytes, url: str = "") -> dict:
    try:
        text = data.decode("utf-8", errors="replace").strip()
    except Exception as e:
        return _result(url=url, content_type="text", error=str(e))
    if not text:
        return _result(url=url, content_type="text", error="empty file")
    return _result(text, content_type="text", url=url)


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_BYTES_EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "xlsx": extract_xlsx,
    "csv": extract_csv,
    "image": extract_image,
    "text": extract_text,
}


def extract_from_bytes(kind: str, data: bytes, url: str = "") -> dict:
    """Dispatch to the right bytes-based extractor. `kind` should come from
    detect_content_kind(). HTML is handled separately via extract_html
    since it takes decoded text, not raw bytes; video/audio are handled
    separately since they take a URL, not bytes the caller fetched."""
    fn = _BYTES_EXTRACTORS.get(kind)
    if fn is None:
        return _result(url=url, content_type=kind or "unknown",
                        error=f"no extractor registered for content kind {kind!r}")
    return fn(data, url)
